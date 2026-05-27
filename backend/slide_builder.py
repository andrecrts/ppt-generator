"""
slide_builder.py
Builds a .pptx presentation from classified slide content using the EDU Template.

Template layout indices (from inspection):
  0  → "Title Slide"           CENTER_TITLE (idx=0) + SUBTITLE (idx=1)
  2  → "1_Title and Content"   TITLE (idx=0) + OBJECT/content (idx=1)
  3  → "2_Title and Content"   TITLE (idx=0) + OBJECT/content (idx=1)
  5  → "Section Header"        TITLE (idx=0) + BODY (idx=1)

UPANA formatting rules applied:
  - Body font: 16pt (default), 14pt if >6 bullets
  - normAutofit enabled on all text frames
"""

import io
from io import BytesIO
from pathlib import Path
from typing import List, Optional

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from classifier import Slide

TEMPLATE_PATH = Path(__file__).parent / "templates" / "EDU Template.pptx"

# Layout indices
_LAYOUT_TITLE = 0         # "Title Slide"
_LAYOUT_CONTENT = 2       # "1_Title and Content"
_LAYOUT_SECTION = 5       # "Section Header"

# Font sizes (UPANA rules)
_FONT_BODY = Pt(16)
_FONT_BODY_COMPACT = Pt(14)
_COMPACT_THRESHOLD = 6    # use smaller font when bullets > this


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _enable_norm_autofit(text_frame) -> None:
    """Replace any existing autofit element with <a:normAutofit/>."""
    txBody = text_frame._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        return

    # Remove conflicting autofit siblings
    for tag in (qn("a:normAutofit"), qn("a:spAutoFit"), qn("a:noAutofit")):
        for el in bodyPr.findall(tag):
            bodyPr.remove(el)

    etree.SubElement(bodyPr, qn("a:normAutofit"))


def _remove_all_slides(prs: Presentation) -> None:
    """
    Remove every existing slide from the presentation, keeping slide layouts
    and masters intact so we can still add new slides from the template.
    """
    slide_id_list = prs.slides._sldIdLst
    NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    for sId in list(slide_id_list):
        r_id = sId.get(f"{{{NS}}}id")
        slide_id_list.remove(sId)
        try:
            prs.part.drop_rel(r_id)
        except Exception:
            pass  # Relationship already gone — safe to ignore


def _get_layout(prs: Presentation, idx: int):
    return prs.slide_masters[0].slide_layouts[idx]


# ---------------------------------------------------------------------------
# Slide-type builders
# ---------------------------------------------------------------------------

def _add_title_slide(prs: Presentation, title: str, subtitle: str = "", author: str = "") -> None:
    slide = prs.slides.add_slide(_get_layout(prs, _LAYOUT_TITLE))
    # Build subtitle: "Subtitle | UPANA" on first line, author on second (if provided)
    # e.g.  subtitle="Semana 2"  → "Semana 2 | UPANA"
    #        subtitle=""          → "UPANA"
    upana_line = f"{subtitle} | UPANA" if subtitle.strip() else "UPANA"
    parts = [upana_line] + ([author.strip()] if author.strip() else [])
    combined_subtitle = "\n".join(parts)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
        elif ph.placeholder_format.idx == 1:
            ph.text = combined_subtitle


_IMG_LEFT_SECTION = 5.5   # inches — image starts here on section slides

def _add_section_slide(prs: Presentation, title: str, img_bytes: Optional[bytes] = None) -> None:
    slide = prs.slides.add_slide(_get_layout(prs, _LAYOUT_SECTION))
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            if img_bytes:
                # Constrain title to the left column so it never overlaps the photo
                ph.left   = Inches(0.5)
                ph.top    = Inches(1.8)
                ph.width  = Inches(_IMG_LEFT_SECTION - 0.9)   # ≈ 4.6" — ends before image
                ph.height = Inches(4.0)
            ph.text = title

    if img_bytes:
        pic = slide.shapes.add_picture(
            BytesIO(img_bytes),
            left=Inches(_IMG_LEFT_SECTION),
            top=Inches(0.0),
            width=Inches(4.5),
            height=Inches(7.5),   # full slide height for a dramatic bleed
        )
        # Move image to the back of the z-order (behind title text)
        sp = pic._element
        spTree = slide.shapes._spTree
        spTree.remove(sp)
        spTree.insert(2, sp)  # 2 = behind other shapes but above background


def _add_content_slide(prs: Presentation, title: str, bullets: List[str], img_bytes: Optional[bytes] = None) -> None:
    slide = prs.slides.add_slide(_get_layout(prs, _LAYOUT_CONTENT))
    font_size = _FONT_BODY_COMPACT if len(bullets) > _COMPACT_THRESHOLD else _FONT_BODY

    IMG_LEFT = 6.4  # inches

    # Narrow the content placeholder to avoid overlap with the image
    if img_bytes:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.left   = Inches(0.688)
                ph.top    = Inches(2.074)
                ph.width  = Inches(IMG_LEFT - 0.8)  # = 5.6"
                ph.height = Inches(4.832)

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = title
        elif idx == 1 and bullets:
            tf = ph.text_frame
            tf.clear()
            _enable_norm_autofit(tf)

            for i, bullet_text in enumerate(bullets):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = bullet_text
                para.level = 0

                for run in para.runs:
                    run.font.size = font_size

    # Add the picture AFTER setting placeholder dimensions
    if img_bytes:
        slide.shapes.add_picture(
            BytesIO(img_bytes),
            left=Inches(IMG_LEFT),
            top=Inches(2.5),
            width=Inches(3.3),
            height=Inches(3.5),
        )


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def build_presentation(
    slides: List[Slide],
    course_name: str = "",
    author_name: str = "",
    images: dict = None,
    template_path: Optional[Path] = None,
) -> bytes:
    """
    Build a .pptx from a list of Slide objects.

    Args:
        slides:        Classified slide list from classifier.py.
        course_name:   Optional subtitle shown on the title slide.
        author_name:   Optional author / teacher name.
        images:        Optional dict mapping slide index (0-based) to image bytes.
        template_path: Optional path to a custom .pptx template.
                       Falls back to the bundled EDU Template when not provided.
                       The template must have compatible layouts at indices 0, 2, 5
                       (Title Slide / Content / Section Header).

    Returns:
        Raw bytes of the generated .pptx file.
    """
    prs = Presentation(str(template_path or TEMPLATE_PATH))
    _remove_all_slides(prs)

    images = images or {}

    for i, slide in enumerate(slides):
        img_bytes = images.get(i)
        if slide.slide_type == "title":
            _add_title_slide(prs, slide.title, subtitle=course_name, author=author_name)
        elif slide.slide_type == "section_header":
            _add_section_slide(prs, slide.title, img_bytes=img_bytes)
        else:
            # "content" or any unknown type → content slide
            _add_content_slide(prs, slide.title, slide.bullets, img_bytes=img_bytes)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
